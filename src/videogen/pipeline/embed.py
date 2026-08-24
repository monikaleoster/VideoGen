"""Real embed step: insert each slide's audio clip into the deck, set to
autoplay on slide entry.

Per specs/2026-08-23-embed-audio-real/requirements.md: a slide with no
audio (None in tts's output) is left untouched, the result is written as
a new file (the source .pptx is never modified), and any failure
propagates rather than being retried — the human re-runs manually via
the existing Reject action.

python-pptx has no dedicated "add audio" method, only `add_movie`, which
works for audio (mime_type="audio/mpeg") but only wires click-to-play
timing. Real autoplay-on-entry requires a direct XML edit: `add_movie`
gives the media node's start condition `<p:cond delay="indefinite"/>`
(wait for the click action), which is changed to `<p:cond delay="0"/>`
(start immediately when the slide's timing root begins) — see
requirements.md's "Technical constraint" section for how this was
determined.
"""

import asyncio
import subprocess
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from videogen.pipeline import workdir
from videogen.pipeline.base import StepState, StepStatus

logger = logging.getLogger(__name__)

# A 1x1 EMU placeholder — this is a narration track, not a video, so its
# on-slide visual footprint isn't a design concern for this phase.
_MEDIA_SIZE = Emu(1)

_TIMING_NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}


@dataclass
class EmbedInput:
    local_pptx_path: str
    audio_paths: list[str | None]


@dataclass
class EmbedOutput:
    updated_pptx_path: str
    slides_embedded: list[bool]


# Module-level state so the CLI and HTTP routes reach the same in-flight run.
state: StepState[EmbedOutput] = StepState()


def _make_autoplay(slide) -> None:
    """Rewrite the just-inserted media's start condition from "wait for
    click" (`delay="indefinite"`) to "start immediately" (`delay="0"`)."""
    cond = slide._element.find(".//p:timing//p:cond", _TIMING_NS)
    if cond is None or cond.get("delay") != "indefinite":
        raise RuntimeError(
            "Expected a click-triggered media timing condition to rewrite for "
            "autoplay, but none was found — python-pptx's add_movie output may "
            "have changed shape."
        )
    cond.set("delay", "0")


def _embed_audio(pptx_path: str, audio_paths: list[str | None], out_path: Path) -> list[bool]:
    """Insert each non-None audio path into its slide, autoplaying on
    entry. Raises on any failure (missing/corrupt audio file, unexpected
    XML shape) — no partial output is written in that case."""
    presentation = Presentation(pptx_path)
    slides_embedded: list[bool] = []

    for i, (slide, audio_path) in enumerate(zip(presentation.slides, audio_paths, strict=True), start=1):
        if audio_path is None:
            logger.debug("Slide %d: no audio, leaving untouched", i)
            slides_embedded.append(False)
            continue

        logger.debug("Slide %d: embedding %s, setting autoplay", i, audio_path)
        slide.shapes.add_movie(
            audio_path,
            left=Emu(0),
            top=Emu(0),
            width=_MEDIA_SIZE,
            height=_MEDIA_SIZE,
            mime_type="audio/mpeg",
        )
        _make_autoplay(slide)
        slides_embedded.append(True)

    presentation.save(str(out_path))
    logger.info("Embedded audio into %d/%d slide(s), saved to %s", sum(slides_embedded), len(slides_embedded), out_path)
    return slides_embedded


async def run_embed(step_input: EmbedInput) -> EmbedOutput:
    """Run real audio embedding, blocking on `state.approval_event` until approved."""
    state.status = StepStatus.RUNNING
    state.output = None
    state.approval_event.clear()
    logger.info("embed starting: local_pptx_path=%s", step_input.local_pptx_path)

    source_path = Path(step_input.local_pptx_path)
    work_dir = Path(tempfile.mkdtemp(prefix="videogen_embed_"))
    out_path = work_dir / f"{source_path.stem}_with_audio.pptx"

    # python-pptx/lxml work is blocking — run it off the event loop thread
    # so the WebSocket status push keeps working during embedding.
    slides_embedded = await asyncio.to_thread(
        _embed_audio, str(source_path), step_input.audio_paths, out_path
    )

    output = EmbedOutput(updated_pptx_path=str(out_path), slides_embedded=slides_embedded)
    state.output = output
    state.status = StepStatus.WAITING_APPROVAL

    await state.approval_event.wait()

    state.status = StepStatus.DONE
    return output


def _embed_all(step_input: EmbedInput) -> EmbedOutput:
    work_dir = workdir.make_work_dir(prefix="videogen_embed_")
    silence_path = work_dir / "silence.mp3"
    _generate_silence(silence_path)

    prs = Presentation(step_input.local_pptx_path)

    slides_embedded: list[bool] = []
    used_placeholder: list[bool] = []
    for slide, audio_path in zip(prs.slides, step_input.audio_paths, strict=True):
        is_placeholder = audio_path is None
        clip_path = str(silence_path) if is_placeholder else audio_path
        _embed_clip(slide, clip_path)
        slides_embedded.append(True)
        used_placeholder.append(is_placeholder)

    base_path = step_input.local_pptx_path.removesuffix(".pptx")
    updated_pptx_path = f"{base_path}_with_audio.pptx"
    prs.save(updated_pptx_path)

    return EmbedOutput(
        updated_pptx_path=updated_pptx_path,
        slides_embedded=slides_embedded,
        used_placeholder=used_placeholder,
    )
