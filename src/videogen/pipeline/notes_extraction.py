"""Real notes-extraction step: local .pptx path -> per-slide speaker notes.

Per specs/2026-08-23-notes-extraction-real/requirements.md: pulls each
slide's raw speaker-notes text via python-pptx, in deck order, stripping
only leading/trailing whitespace (no other normalization). A slide with
no notes slide, or an empty/whitespace-only one, is flagged via
`has_notes` but still produces a normal result and still gates on human
approval like any other slide.

Per specs/2026-08-23-notes-text-files/requirements.md: each slide's
stripped notes text is also written to its own `.txt` file, one per
slide, always — including an empty file for a `has_notes=False` slide —
in a work directory created via `workdir.make_work_dir`.
"""

import asyncio
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation

from videogen.pipeline import workdir
from videogen.pipeline.base import StepState, StepStatus


@dataclass
class NotesExtractionInput:
    local_pptx_path: str


@dataclass
class NotesExtractionOutput:
    slide_count: int
    notes: list[str]
    has_notes: list[bool]
    notes_file_paths: list[str]


# Module-level state so the CLI and HTTP routes reach the same in-flight run.
state: StepState[NotesExtractionOutput] = StepState()


def _extract_notes(pptx_path: str) -> tuple[list[str], list[bool]]:
    """Read each slide's speaker notes, in deck order, via python-pptx.

    A missing notes slide and an empty/whitespace-only one are both
    represented the same way: `""` in `notes`, `False` in `has_notes`.
    """
    presentation = Presentation(pptx_path)

    notes: list[str] = []
    has_notes: list[bool] = []
    for slide in presentation.slides:
        raw_text = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        stripped = raw_text.strip()
        notes.append(stripped)
        has_notes.append(bool(stripped))

    return notes, has_notes


def _write_notes_files(notes: list[str], work_dir: Path) -> list[str]:
    """Write each slide's stripped notes text to its own file, one per
    slide, always — including an empty file for a no-notes slide."""
    notes_file_paths: list[str] = []
    for i, text in enumerate(notes, start=1):
        out_path = work_dir / f"slide_{i:02d}_notes.txt"
        out_path.write_text(text)
        notes_file_paths.append(str(out_path))
    return notes_file_paths


async def run_notes_extraction(step_input: NotesExtractionInput) -> NotesExtractionOutput:
    """Run the real extraction, blocking on `state.approval_event` until approved."""
    state.status = StepStatus.RUNNING
    state.output = None
    state.approval_event.clear()

    # python-pptx parsing is blocking (file I/O + XML parsing) — run it off
    # the event loop thread so the WebSocket status push keeps working.
    notes, has_notes = await asyncio.to_thread(_extract_notes, step_input.local_pptx_path)

    work_dir = workdir.make_work_dir(prefix="videogen_notes_")
    notes_file_paths = await asyncio.to_thread(_write_notes_files, notes, work_dir)

    output = NotesExtractionOutput(
        slide_count=len(notes),
        notes=notes,
        has_notes=has_notes,
        notes_file_paths=notes_file_paths,
    )
    state.output = output
    state.status = StepStatus.WAITING_APPROVAL

    await state.approval_event.wait()

    state.status = StepStatus.DONE
    return output
