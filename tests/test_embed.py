import asyncio
import subprocess
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from videogen.pipeline import workdir
from videogen.pipeline.base import StepStatus
from videogen.pipeline.embed import EmbedInput, run_embed, state


@pytest.fixture(autouse=True)
def reset_state():
    state.status = StepStatus.PENDING
    state.output = None
    state.approval_event = asyncio.Event()
    workdir.set_tmp_root(None)
    yield
    state.status = StepStatus.PENDING
    state.output = None
    state.approval_event = asyncio.Event()
    workdir.set_tmp_root(None)


@pytest.fixture
def deck_path(tmp_path: Path) -> str:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.slides.add_slide(prs.slide_layouts[6])
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return str(path)


@pytest.fixture
def real_clip(tmp_path: Path) -> str:
    """A real, short MP3 clip (not silence) to distinguish from the
    placeholder in tests."""
    out_path = tmp_path / "clip.mp3"
    subprocess.run(
        [
            "ffmpeg",
            "-y",
            "-f",
            "lavfi",
            "-i",
            "sine=frequency=440:duration=2",
            "-q:a",
            "9",
            "-acodec",
            "libmp3lame",
            str(out_path),
        ],
        check=True,
        capture_output=True,
    )
    return str(out_path)


def _audio_shapes(slide):
    return [shape for shape in slide.shapes if shape._element.find(f".//{qn('a:audioFile')}") is not None]


def _has_autoplay_timing(slide) -> bool:
    timing = slide._element.find(qn("p:timing"))
    return timing is not None and timing.find(f".//{qn('p:cmd')}") is not None


async def test_step_blocks_until_approved(deck_path: str, real_clip: str) -> None:
    task = asyncio.create_task(run_embed(EmbedInput(local_pptx_path=deck_path, audio_paths=[real_clip, None, None])))

    await asyncio.wait_for(_wait_for_status(StepStatus.WAITING_APPROVAL), timeout=5.0)
    assert state.status == StepStatus.WAITING_APPROVAL

    await asyncio.sleep(0.05)
    assert not task.done()

    task.cancel()
    try:
        await task
    except asyncio.CancelledError:
        pass


async def test_step_resumes_and_completes_after_approval(deck_path: str, real_clip: str) -> None:
    task = asyncio.create_task(run_embed(EmbedInput(local_pptx_path=deck_path, audio_paths=[real_clip, None, None])))

    await asyncio.wait_for(_wait_for_status(StepStatus.WAITING_APPROVAL), timeout=5.0)
    state.approval_event.set()
    output = await asyncio.wait_for(task, timeout=5.0)

    assert state.status == StepStatus.DONE
    assert output.updated_pptx_path
    assert output.slides_embedded == [True, True, True]
    assert output.used_placeholder == [False, True, True]
    assert state.output is output


async def test_real_clip_and_placeholder_are_both_embedded_with_autoplay(deck_path: str, real_clip: str) -> None:
    task = asyncio.create_task(run_embed(EmbedInput(local_pptx_path=deck_path, audio_paths=[real_clip, None, None])))
    await asyncio.wait_for(_wait_for_status(StepStatus.WAITING_APPROVAL), timeout=5.0)
    state.approval_event.set()
    output = await asyncio.wait_for(task, timeout=5.0)

    prs = Presentation(output.updated_pptx_path)
    assert len(prs.slides) == 3

    for slide in prs.slides:
        audio_shapes = _audio_shapes(slide)
        assert len(audio_shapes) == 1, "each slide must have exactly one audio element"
        assert _has_autoplay_timing(slide), "audio must be set to autoplay on slide entry"

        shape = audio_shapes[0]
        assert shape.width == shape.height == 1, "media placeholder/icon must be hidden"


async def test_rerunning_the_step_does_not_duplicate_audio(deck_path: str, real_clip: str) -> None:
    task = asyncio.create_task(run_embed(EmbedInput(local_pptx_path=deck_path, audio_paths=[real_clip, None, None])))
    await asyncio.wait_for(_wait_for_status(StepStatus.WAITING_APPROVAL), timeout=5.0)
    state.approval_event.set()
    await asyncio.wait_for(task, timeout=5.0)

    state.status = StepStatus.PENDING
    state.output = None
    state.approval_event = asyncio.Event()

    # Re-run against the same pristine source deck, as the real pipeline
    # does after a per-slide TTS regeneration (runner.py always passes the
    # download step's untouched local_pptx_path, never embed's own output).
    task2 = asyncio.create_task(run_embed(EmbedInput(local_pptx_path=deck_path, audio_paths=[real_clip, None, None])))
    await asyncio.wait_for(_wait_for_status(StepStatus.WAITING_APPROVAL), timeout=5.0)
    state.approval_event.set()
    output2 = await asyncio.wait_for(task2, timeout=5.0)

    prs = Presentation(output2.updated_pptx_path)
    for slide in prs.slides:
        assert len(_audio_shapes(slide)) == 1


async def test_work_dir_nests_under_shared_tmp_root_when_set(
    tmp_path: Path, deck_path: str, real_clip: str
) -> None:
    custom_root = tmp_path / "shared_root"
    workdir.set_tmp_root(str(custom_root))

    task = asyncio.create_task(run_embed(EmbedInput(local_pptx_path=deck_path, audio_paths=[real_clip, None, None])))
    await asyncio.wait_for(_wait_for_status(StepStatus.WAITING_APPROVAL), timeout=5.0)
    state.approval_event.set()
    await asyncio.wait_for(task, timeout=5.0)

    # embed's own work dir (holding the generated silence.mp3) isn't
    # surfaced in its output, so confirm nesting by inspecting the root
    # directly: exactly one videogen_embed_* dir was created under it.
    embed_dirs = list(custom_root.glob("videogen_embed_*"))
    assert len(embed_dirs) == 1
    assert (embed_dirs[0] / "silence.mp3").exists()


async def _wait_for_status(target: StepStatus) -> None:
    while state.status != target:
        await asyncio.sleep(0.01)
